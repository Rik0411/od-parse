"""
PowerPoint Presentation (.pptx) Pipeline for parsing Microsoft PowerPoint files.

Extracts text, tables, and images from PowerPoint presentations, organized by slides.
Saves extracted images to disk and returns their paths.
"""

import os
from pathlib import Path
from typing import Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

from od_parse.utils.logging_utils import get_logger

logger = get_logger(__name__)

# Import image description function
try:
    from od_parse.excel.gemini_client import generate_excel_image_description
except ImportError:
    generate_excel_image_description = None


def runPptxPipeline(file_path: str, args, output_dir: Optional[Path] = None) -> Dict[str, Any]:
    """
    Parses a .pptx file for slides, text, tables, and images.
    Saves extracted images to disk and returns their paths.
    
    Args:
        file_path: Path to the PowerPoint presentation file
        args: argparse.Namespace object with .mech attribute
        output_dir: Optional output directory (if None, creates output/{filename}_images/)
        
    Returns:
        Dictionary with structure:
        {
            "file_type": "pptx",
            "slides": [
                {
                    "slide_number": 1,
                    "title": "...",
                    "text_content": [...],
                    "tables": [...],
                    "images": [{"path": "...", "format": "...", "description": "..."}]
                }
            ]
        }
    """
    logger.info(f"Starting PPTX Pipeline for: {file_path}")
    
    file_path_obj = Path(file_path)
    if not file_path_obj.exists():
        raise FileNotFoundError(f"PowerPoint presentation not found: {file_path}")
    
    prs = Presentation(file_path)
    
    # Create an output directory for images based on the input filename
    if output_dir is None:
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_dir = Path("output") / f"{base_name}_images"
    else:
        output_dir = Path(output_dir)
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Get API key from environment if --mech flag is set
    api_key = os.getenv("GOOGLE_API_KEY") if args.mech else None
    if args.mech and not api_key:
        logger.warning("--mech flag is set but GOOGLE_API_KEY not found. Images will not be captioned.")
    
    result = {
        "file_type": "pptx",
        "slides": []
    }

    image_count = 0

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "title": "",
            "text_content": [],
            "tables": [],
            "images": []
        }

        # Extract Title (if available)
        try:
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text
        except:
            pass  # Title might not exist

        # Iterate all shapes
        for shape in slide.shapes:
            # 1. Text extraction
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_data["text_content"].append(text)

            # 2. Table extraction
            if shape.has_table:
                table_data = []
                keys = None
                for r_idx, row in enumerate(shape.table.rows):
                    text = (cell.text_frame.text.strip() for cell in row.cells)
                    if r_idx == 0:
                        keys = tuple(text)
                        continue
                    # Handle varying row lengths
                    row_values = list(text)
                    if keys and len(row_values) == len(keys):
                        row_data = dict(zip(keys, row_values))
                        table_data.append(row_data)
                    else:
                        table_data.append(row_values)
                slide_data["tables"].append(table_data)

            # 3. Image extraction
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    ext = shape.image.ext
                    
                    # Generate unique filename
                    image_filename = f"slide_{i+1}_image_{image_count}.{ext}"
                    image_path = output_dir / image_filename
                    
                    # Save to disk (ALWAYS happen, regardless of flags)
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    image_info = {
                        "path": str(image_path),
                        "format": ext,
                        "description": "Image extracted from slide"
                    }
                    
                    # Semantic Analysis (ONLY if --mech is on)
                    if args.mech and generate_excel_image_description and api_key:
                        try:
                            # Convert bytes to PIL Image (required by generate_excel_image_description)
                            pil_image = Image.open(io.BytesIO(image_bytes))
                            if pil_image.mode != 'RGB':
                                pil_image = pil_image.convert('RGB')
                            
                            desc = generate_excel_image_description(pil_image, api_key)
                            image_info["description"] = desc
                        except Exception as e:
                            logger.warning(f"Failed to generate description for image on slide {i+1}: {e}")
                    
                    slide_data["images"].append(image_info)
                    image_count += 1
                    logger.info(f"Extracted and saved image from slide {i+1} to: {image_path}")
                except Exception as e:
                    logger.warning(f"Failed to extract image from slide {i+1}: {e}")

        result["slides"].append(slide_data)

    logger.info(f"PPTX extraction complete. Parsed {len(result['slides'])} slides.")
    return result
